"""
╔══════════════════════════════════════════════════════════════════════════════╗
║                    NEURAL COUNCIL - THE ABSOLUTE PINNACLE                    ║
╚══════════════════════════════════════════════════════════════════════════════╝

The most beautiful AI interface. 4 models working as one.
Neon pink/sunset aesthetic. Glassmorphism. Premium animations.
"""

import streamlit as st
import os
import re
import base64
import json  # Needed for Jupyter notebook parsing
from datetime import datetime
from dotenv import load_dotenv

load_dotenv()

st.set_page_config(
    page_title="Neural Council",
    page_icon="🧠",
    layout="wide",
    initial_sidebar_state="expanded"
)

import council


def render_message_with_images(content: str):
    """Render message and display any embedded images."""
    st.markdown(content)
    
    # Extract and display images
    patterns = [
        r'!\[[^\]]*\]\(([^\)]+)\)',
        r'(https://[^\s<>"]+\.(?:png|jpg|jpeg|gif|webp))',
        r'(https://oaidalleapiprodscus[^\s<>"]+)',
    ]
    
    displayed = set()
    for pattern in patterns:
        for match in re.finditer(pattern, content):
            url = match.group(1) if match.lastindex else match.group(0)
            if url and url not in displayed:
                try:
                    st.image(url, width=450)
                    displayed.add(url)
                except:
                    pass


def smart_display_content(content: str, max_preview: int = 500, is_user: bool = False):
    """
    SMART DISPLAY: Collapse large content to prevent screen flooding.
    - Shows preview for user messages with files
    - Collapses code blocks
    - Handles long responses gracefully
    """
    if not content:
        return
    
    # Detect if content has embedded files/code
    has_file = "[FILE:" in content or "[EXCEL:" in content or "[CSV:" in content or "[WORD:" in content
    has_code = "```" in content and len(content) > 2000
    
    # For USER messages with files - show clean summary
    if is_user and has_file:
        # Extract just the query part (before file content)
        query_part = content.split("[FILE:")[0].split("[EXCEL:")[0].split("[CSV:")[0].split("[WORD:")[0].strip()
        
        # Count files
        file_count = content.count("[FILE:") + content.count("[EXCEL:") + content.count("[CSV:") + content.count("[WORD:") + content.count("[JUPYTER:") + content.count("[POWERPOINT:")
        
        # Show clean preview
        st.markdown(query_part[:max_preview] + ("..." if len(query_part) > max_preview else ""))
        st.caption(f"📎 {file_count} file(s) attached")
        
        # Collapsible file details
        with st.expander("📂 View attached file content", expanded=False):
            st.markdown(content)
        return
    
    # For ASSISTANT messages with very long code blocks
    if has_code and len(content) > 3000:
        # Split into text and code parts
        parts = content.split("```")
        
        # Show text preview
        preview_text = parts[0][:max_preview]
        if len(parts[0]) > max_preview:
            preview_text += "..."
        st.markdown(preview_text)
        
        # Show code in expander
        if len(parts) >= 3:
            code_content = parts[1]
            lang = code_content.split("\n")[0] if "\n" in code_content else ""
            code = "\n".join(code_content.split("\n")[1:]) if "\n" in code_content else code_content
            
            with st.expander(f"📜 View code ({len(code)} chars)", expanded=True):
                st.code(code, language=lang if lang else "text")
            
            # Show any remaining text
            if len(parts) > 2:
                remaining = "".join(parts[2:])[:500]
                if remaining.strip():
                    st.markdown(remaining)
        return
    
    # For normal content - just render with images
    render_message_with_images(content)


# ═══════════════════════════════════════════════════════════════════════════════
# STUNNING CSS - NEON PINK/SUNSET AESTHETIC
# ═══════════════════════════════════════════════════════════════════════════════

NEON_CSS = """
<style>
@import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700&family=Space+Grotesk:wght@400;500;600;700&display=swap');

:root {
    --bg-dark: #0a0a0f;
    --bg-card: rgba(20, 15, 30, 0.8);
    --neon-pink: #ff1493;
    --neon-magenta: #ff00ff;
    --sunset-orange: #ff6b35;
    --sunset-pink: #ff8577;
    --soft-pink: #ffb6c1;
    --text-primary: #ffffff;
    --text-secondary: rgba(255, 255, 255, 0.7);
    --glass-border: rgba(255, 20, 147, 0.3);
}

* {
    font-family: 'Inter', -apple-system, BlinkMacSystemFont, sans-serif;
}

.stApp {
    background: linear-gradient(135deg, 
        #0a0a0f 0%, 
        #1a0a1f 25%,
        #150520 50%,
        #0f0515 75%,
        #0a0a0f 100%);
    background-attachment: fixed;
}

/* Animated background glow */
.stApp::before {
    content: '';
    position: fixed;
    top: 0;
    left: 0;
    right: 0;
    bottom: 0;
    background: 
        radial-gradient(ellipse at 20% 20%, rgba(255, 20, 147, 0.15) 0%, transparent 50%),
        radial-gradient(ellipse at 80% 80%, rgba(255, 107, 53, 0.1) 0%, transparent 50%),
        radial-gradient(ellipse at 50% 50%, rgba(255, 0, 255, 0.05) 0%, transparent 70%);
    pointer-events: none;
    z-index: 0;
}

/* Sidebar - Glassmorphism */
[data-testid="stSidebar"] {
    background: linear-gradient(180deg, 
        rgba(20, 10, 35, 0.95) 0%, 
        rgba(15, 5, 25, 0.98) 100%);
    border-right: 1px solid var(--glass-border);
    backdrop-filter: blur(20px);
}

[data-testid="stSidebar"]::before {
    content: '';
    position: absolute;
    top: 0;
    left: 0;
    right: 0;
    height: 3px;
    background: linear-gradient(90deg, var(--neon-pink), var(--neon-magenta), var(--sunset-orange));
}

/* Buttons - Neon glow effect */
.stButton > button {
    background: linear-gradient(135deg, 
        rgba(255, 20, 147, 0.2) 0%, 
        rgba(255, 0, 255, 0.15) 100%);
    color: var(--text-primary);
    border: 1px solid var(--glass-border);
    border-radius: 12px;
    font-weight: 600;
    padding: 12px 24px;
    min-height: 48px;
    backdrop-filter: blur(10px);
    transition: all 0.3s cubic-bezier(0.4, 0, 0.2, 1);
    text-shadow: 0 0 10px rgba(255, 20, 147, 0.5);
}

.stButton > button:hover {
    background: linear-gradient(135deg, 
        rgba(255, 20, 147, 0.4) 0%, 
        rgba(255, 0, 255, 0.3) 100%);
    border-color: var(--neon-pink);
    transform: translateY(-2px);
    box-shadow: 
        0 0 20px rgba(255, 20, 147, 0.4),
        0 0 40px rgba(255, 20, 147, 0.2),
        0 10px 30px rgba(0, 0, 0, 0.3);
}

/* Input fields - Glass effect */
.stTextInput > div > div > input,
.stTextArea > div > div > textarea,
[data-testid="stChatInput"] textarea {
    background: rgba(20, 10, 35, 0.6) !important;
    color: var(--text-primary) !important;
    border: 1px solid var(--glass-border) !important;
    border-radius: 12px !important;
    backdrop-filter: blur(10px);
    transition: all 0.3s ease;
}

.stTextInput > div > div > input:focus,
.stTextArea > div > div > textarea:focus,
[data-testid="stChatInput"] textarea:focus {
    border-color: var(--neon-pink) !important;
    box-shadow: 0 0 20px rgba(255, 20, 147, 0.3) !important;
}

/* Chat messages - Premium glass cards */
[data-testid="stChatMessage"] {
    background: linear-gradient(135deg, 
        rgba(25, 15, 40, 0.7) 0%, 
        rgba(20, 10, 30, 0.8) 100%) !important;
    border: 1px solid rgba(255, 20, 147, 0.2);
    border-radius: 20px;
    padding: 20px;
    margin: 12px 0;
    backdrop-filter: blur(15px);
    transition: all 0.3s ease;
}

[data-testid="stChatMessage"]:hover {
    border-color: rgba(255, 20, 147, 0.4);
    box-shadow: 0 0 30px rgba(255, 20, 147, 0.15);
}

/* Agent badges - Glowing pills */
.agent-badge {
    display: inline-block;
    padding: 6px 16px;
    border-radius: 30px;
    font-size: 0.85em;
    font-weight: 600;
    margin-bottom: 10px;
    letter-spacing: 0.5px;
}

.agent-emperor {
    background: linear-gradient(90deg, #ffd700, #ff8c00);
    color: #000;
    box-shadow: 0 0 20px rgba(255, 215, 0, 0.5);
}

.agent-strategist {
    background: linear-gradient(90deg, #00bfff, #1e90ff);
    color: #fff;
    box-shadow: 0 0 20px rgba(0, 191, 255, 0.5);
}

.agent-executor {
    background: linear-gradient(90deg, #00ff88, #00cc66);
    color: #000;
    box-shadow: 0 0 20px rgba(0, 255, 136, 0.5);
}

.agent-sage {
    background: linear-gradient(90deg, #bf00ff, #8000ff);
    color: #fff;
    box-shadow: 0 0 20px rgba(191, 0, 255, 0.5);
}

/* Headers - Gradient text */
h1, h2, h3 {
    background: linear-gradient(90deg, var(--neon-pink), var(--sunset-orange));
    -webkit-background-clip: text;
    -webkit-text-fill-color: transparent;
    background-clip: text;
    font-family: 'Space Grotesk', sans-serif;
    font-weight: 700;
}

/* Code blocks - Dark glass */
pre, code {
    background: rgba(10, 5, 20, 0.8) !important;
    border: 1px solid rgba(255, 20, 147, 0.2) !important;
    border-radius: 12px !important;
}

/* Expander - Glass style */
.streamlit-expanderHeader {
    background: rgba(20, 10, 35, 0.6) !important;
    border: 1px solid var(--glass-border) !important;
    border-radius: 12px !important;
}

/* Selectbox - Glass */
.stSelectbox > div > div {
    background: rgba(20, 10, 35, 0.6) !important;
    border: 1px solid var(--glass-border) !important;
    border-radius: 12px !important;
}

/* File uploader */
[data-testid="stFileUploader"] {
    background: rgba(20, 10, 35, 0.5);
    border: 2px dashed var(--glass-border);
    border-radius: 16px;
    padding: 20px;
}

/* Divider - Gradient */
hr {
    border: none;
    height: 1px;
    background: linear-gradient(90deg, 
        transparent 0%, 
        var(--neon-pink) 50%, 
        transparent 100%);
    margin: 20px 0;
}

/* Status indicator - Pulsing */
.stStatus {
    background: rgba(20, 10, 35, 0.8) !important;
    border: 1px solid var(--glass-border) !important;
    border-radius: 16px !important;
}

/* Scrollbar - Neon */
::-webkit-scrollbar {
    width: 8px;
    height: 8px;
}

::-webkit-scrollbar-track {
    background: rgba(10, 5, 20, 0.5);
}

::-webkit-scrollbar-thumb {
    background: linear-gradient(180deg, var(--neon-pink), var(--neon-magenta));
    border-radius: 10px;
}

/* Mobile responsive */
@media (max-width: 768px) {
    .stApp { padding: 0.5rem !important; }
    h1 { font-size: 1.8em !important; }
    h2 { font-size: 1.4em !important; }
    .stButton > button { width: 100% !important; }
    [data-testid="column"] { width: 100% !important; flex: 1 1 100% !important; }
}

/* Animations */
@keyframes glow-pulse {
    0%, 100% { box-shadow: 0 0 20px rgba(255, 20, 147, 0.3); }
    50% { box-shadow: 0 0 40px rgba(255, 20, 147, 0.5); }
}

@keyframes float {
    0%, 100% { transform: translateY(0); }
    50% { transform: translateY(-5px); }
}

/* Logo animation */
.logo-container {
    animation: float 3s ease-in-out infinite;
}

/* Tabs styling */
.stTabs [data-baseweb="tab-list"] {
    gap: 8px;
    background: transparent;
}

.stTabs [data-baseweb="tab"] {
    background: rgba(20, 10, 35, 0.5);
    border: 1px solid var(--glass-border);
    border-radius: 12px;
    color: var(--text-secondary);
    padding: 10px 20px;
}

.stTabs [data-baseweb="tab"][aria-selected="true"] {
    background: linear-gradient(135deg, rgba(255, 20, 147, 0.3), rgba(255, 0, 255, 0.2));
    border-color: var(--neon-pink);
    color: var(--text-primary);
}

/* Success/Error messages */
.stSuccess, .stError, .stWarning, .stInfo {
    background: rgba(20, 10, 35, 0.8) !important;
    border-radius: 12px !important;
    border: 1px solid var(--glass-border) !important;
}
</style>
"""

st.markdown(NEON_CSS, unsafe_allow_html=True)

# ═══════════════════════════════════════════════════════════════════════════════
# SCREEN CAPTURE
# ═══════════════════════════════════════════════════════════════════════════════

SCREEN_CAPTURE_JS = """
<script>
async function captureScreen() {
    try {
        const stream = await navigator.mediaDevices.getDisplayMedia({video: {cursor: "always"}, audio: false});
        const track = stream.getVideoTracks()[0];
        const imageCapture = new ImageCapture(track);
        const bitmap = await imageCapture.grabFrame();
        
        const canvas = document.createElement('canvas');
        canvas.width = bitmap.width;
        canvas.height = bitmap.height;
        canvas.getContext('2d').drawImage(bitmap, 0, 0);
        
        const dataUrl = canvas.toDataURL('image/png');
        track.stop();
        
        window.parent.postMessage({type: 'screenshot', data: dataUrl}, '*');
        document.getElementById('screenshot-status').innerText = '✅ Captured!';
        document.getElementById('screenshot-data').value = dataUrl;
    } catch (e) {
        document.getElementById('screenshot-status').innerText = '❌ ' + e.message;
    }
}
</script>
<div style="margin: 10px 0;">
    <button onclick="captureScreen()" style="
        background: linear-gradient(135deg, #ff1493, #ff00ff);
        color: white;
        border: none;
        padding: 12px 24px;
        border-radius: 12px;
        cursor: pointer;
        font-weight: 600;
        font-size: 14px;
        box-shadow: 0 0 20px rgba(255, 20, 147, 0.4);
        transition: all 0.3s ease;
    ">
        📸 Capture Screen
    </button>
    <span id="screenshot-status" style="margin-left: 12px; color: rgba(255,255,255,0.7);"></span>
    <input type="hidden" id="screenshot-data" name="screenshot">
</div>
"""

# ═══════════════════════════════════════════════════════════════════════════════
# SITE PASSWORD GATE - Before login
# ═══════════════════════════════════════════════════════════════════════════════

SITE_PASSWORD = os.getenv("SITE_PASSWORD", "neural2024")  # Change this or set env var

# JavaScript for persistent site unlock - ENHANCED
SITE_UNLOCK_JS = """
<script>
// ENHANCED SITE UNLOCK PERSISTENCE
// Uses: localStorage + sessionStorage + cookie + URL params

function saveSiteUnlock() {
    localStorage.setItem('nc_site_unlocked', 'true');
    sessionStorage.setItem('nc_site_unlocked', 'true');
    // Also set a cookie for extra persistence (7 days)
    document.cookie = 'nc_site_unlocked=true; max-age=604800; path=/; SameSite=Lax';
}

function isSiteUnlocked() {
    // Check all storage methods
    if (localStorage.getItem('nc_site_unlocked') === 'true') return true;
    if (sessionStorage.getItem('nc_site_unlocked') === 'true') return true;
    if (document.cookie.includes('nc_site_unlocked=true')) return true;
    return false;
}

function clearSiteUnlock() {
    localStorage.removeItem('nc_site_unlocked');
    sessionStorage.removeItem('nc_site_unlocked');
    document.cookie = 'nc_site_unlocked=; max-age=0; path=/';
}

// On load: Sync all storage methods and set URL param
(function() {
    const unlocked = isSiteUnlocked();
    if (unlocked) {
        // Sync to all storage methods
        saveSiteUnlock();
        
        // Add to URL params so Streamlit can read it
        const url = new URL(window.location);
        if (!url.searchParams.has('nc_unlocked')) {
            url.searchParams.set('nc_unlocked', 'true');
            window.location.replace(url.toString());
        }
    }
})();
</script>
"""
st.markdown(SITE_UNLOCK_JS, unsafe_allow_html=True)

# Define params early - needed for site unlock persistence
params = st.query_params

# Check URL params for stored unlock status
if params.get("nc_unlocked") == "true":
    st.session_state.site_unlocked = True

# Check if user has passed site gate
if "site_unlocked" not in st.session_state:
    st.session_state.site_unlocked = False

if not st.session_state.site_unlocked:
    st.markdown("""
    <div style="text-align: center; padding: 100px 20px;">
        <div style="font-size: 4em; margin-bottom: 20px;">🔐</div>
        <h1 style="
            font-size: 2.5em;
            background: linear-gradient(90deg, #ff1493, #ff00ff);
            -webkit-background-clip: text;
            -webkit-text-fill-color: transparent;
        ">Access Required</h1>
        <p style="color: rgba(255,255,255,0.5); margin-top: 10px;">Enter the site password to continue</p>
    </div>
    """, unsafe_allow_html=True)
    
    col1, col2, col3 = st.columns([1, 1.5, 1])
    with col2:
        site_pass = st.text_input("Site Password", type="password", key="site_pass_input", placeholder="Enter access code...")
        
        if st.button("🔓 Unlock", use_container_width=True, key="site_unlock_btn"):
            if site_pass == SITE_PASSWORD:
                st.session_state.site_unlocked = True
                # Save to localStorage
                st.markdown('<script>saveSiteUnlock();</script>', unsafe_allow_html=True)
                st.success("✅ Access granted!")
                st.rerun()
            else:
                st.error("❌ Incorrect password")
    
    st.stop()

# ═══════════════════════════════════════════════════════════════════════════════
# AUTHENTICATION WITH PERSISTENT SESSIONS
# ═══════════════════════════════════════════════════════════════════════════════

# params already defined above for site unlock persistence

# JavaScript for persistent session storage - ENHANCED
SESSION_STORAGE_JS = """
<script>
// ENHANCED SESSION PERSISTENCE
// Uses: localStorage + sessionStorage + cookie + URL params

// Save session to ALL storage methods
function saveSession(token, userId, email) {
    // localStorage (persists across browser close)
    localStorage.setItem('nc_token', token);
    localStorage.setItem('nc_user_id', userId);
    localStorage.setItem('nc_email', email);
    
    // sessionStorage (backup for same-tab)
    sessionStorage.setItem('nc_token', token);
    sessionStorage.setItem('nc_user_id', userId);
    sessionStorage.setItem('nc_email', email);
    
    // Cookie (7 day persistence, survives everything)
    document.cookie = `nc_token=${token}; max-age=604800; path=/; SameSite=Lax`;
    document.cookie = `nc_user_id=${userId}; max-age=604800; path=/; SameSite=Lax`;
    document.cookie = `nc_email=${email}; max-age=604800; path=/; SameSite=Lax`;
}

// Get session from ANY available source
function getSession() {
    // Priority: localStorage > sessionStorage > cookie
    let token = localStorage.getItem('nc_token') || sessionStorage.getItem('nc_token');
    let userId = localStorage.getItem('nc_user_id') || sessionStorage.getItem('nc_user_id');
    let email = localStorage.getItem('nc_email') || sessionStorage.getItem('nc_email');
    
    // Fallback to cookies if not in storage
    if (!token) {
        const cookies = document.cookie.split(';').reduce((acc, c) => {
            const [k, v] = c.trim().split('=');
            acc[k] = v;
            return acc;
        }, {});
        token = cookies['nc_token'];
        userId = cookies['nc_user_id'];
        email = cookies['nc_email'];
    }
    
    return { token, userId, email };
}

// Clear session from ALL storage
function clearSession() {
    // localStorage
    localStorage.removeItem('nc_token');
    localStorage.removeItem('nc_user_id');
    localStorage.removeItem('nc_email');
    
    // sessionStorage
    sessionStorage.removeItem('nc_token');
    sessionStorage.removeItem('nc_user_id');
    sessionStorage.removeItem('nc_email');
    
    // Cookies
    document.cookie = 'nc_token=; max-age=0; path=/';
    document.cookie = 'nc_user_id=; max-age=0; path=/';
    document.cookie = 'nc_email=; max-age=0; path=/';
    
    // Clear from URL
    const url = new URL(window.location);
    url.searchParams.delete('nc_token');
    url.searchParams.delete('nc_user_id');
    url.searchParams.delete('nc_email');
    window.history.replaceState({}, '', url);
}

// On page load: Restore session to URL params for Streamlit
(function restoreSession() {
    const session = getSession();
    const url = new URL(window.location);
    
    // If we have a stored session but it's not in URL, add it
    if (session.token && session.userId && !url.searchParams.has('nc_token')) {
        // Re-save to ensure all storage methods are synced
        saveSession(session.token, session.userId, session.email || '');
        
        // Add to URL for Streamlit
        url.searchParams.set('nc_token', session.token);
        url.searchParams.set('nc_user_id', session.userId);
        url.searchParams.set('nc_email', session.email || '');
        window.location.replace(url.toString());
    }
})();

// TAB PERSISTENCE: Restore session when user returns to tab
document.addEventListener('visibilitychange', function() {
    if (document.visibilityState === 'visible') {
        const session = getSession();
        const url = new URL(window.location);
        
        // If session exists but not in URL, restore it
        if (session.token && session.userId && !url.searchParams.has('nc_token')) {
            url.searchParams.set('nc_token', session.token);
            url.searchParams.set('nc_user_id', session.userId);
            url.searchParams.set('nc_email', session.email || '');
            window.location.replace(url.toString());
        }
    }
});
</script>
"""

# Inject session storage JS
st.markdown(SESSION_STORAGE_JS, unsafe_allow_html=True)

# Check for token in query params (for localStorage restoration)
def get_current_user():
    """
    Get current user with FULL PERSISTENCE:
    1. Check session_state (fastest)
    2. Check query params for stored token (survives server restart)
    3. Verify token with Supabase
    """
    # Already in session
    if "user" in st.session_state and st.session_state.user:
        return st.session_state.user
    
    # Check query params for restored token from localStorage
    stored_token = params.get("nc_token")
    stored_user_id = params.get("nc_user_id") 
    stored_email = params.get("nc_email")
    
    if stored_token and stored_user_id:
        # Verify token is still valid with Supabase
        verified_user = council.verify_token(stored_token)
        if verified_user:
            # Token valid - restore session!
            user_data = {
                "id": verified_user.get("id", stored_user_id),
                "email": verified_user.get("email", stored_email),
                "token": stored_token
            }
            st.session_state.user = user_data
            return user_data
        else:
            # Token expired - clear params
            if "nc_token" in st.query_params:
                del st.query_params["nc_token"]
            if "nc_user_id" in st.query_params:
                del st.query_params["nc_user_id"]
            if "nc_email" in st.query_params:
                del st.query_params["nc_email"]
    
    return None

def logout():
    if "user" in st.session_state:
        del st.session_state.user
    if "session_id" in st.session_state:
        del st.session_state.session_id
    # Clear query params
    for key in ["nc_token", "nc_user_id", "nc_email"]:
        if key in st.query_params:
            del st.query_params[key]
    # Clear localStorage via JS
    st.markdown('<script>clearSession();</script>', unsafe_allow_html=True)
    st.rerun()

user = get_current_user()

# ═══════════════════════════════════════════════════════════════════════════════
# LOGIN PAGE - STUNNING
# ═══════════════════════════════════════════════════════════════════════════════

if not user:
    st.markdown("""
    <div style="text-align: center; padding: 60px 20px;">
        <div class="logo-container" style="font-size: 5em; margin-bottom: 20px;">🧠</div>
        <h1 style="
            font-size: 3.5em;
            margin: 0;
            background: linear-gradient(90deg, #ff1493, #ff00ff, #ff6b35);
            -webkit-background-clip: text;
            -webkit-text-fill-color: transparent;
            font-family: 'Space Grotesk', sans-serif;
            font-weight: 700;
            letter-spacing: -2px;
        ">NEURAL COUNCIL</h1>
        <p style="
            color: rgba(255,255,255,0.6);
            font-size: 1.3em;
            margin-top: 15px;
            font-weight: 300;
        ">4 AI Models. One Unified Intelligence.</p>
        <p style="
            color: rgba(255,255,255,0.4);
            font-size: 0.95em;
            margin-top: 8px;
        ">Claude Opus • Claude Sonnet • GPT-5.2 • DeepSeek V3</p>
    </div>
    """, unsafe_allow_html=True)
    
    col1, col2, col3 = st.columns([1, 1.5, 1])
    with col2:
        tab1, tab2 = st.tabs(["🔐 Sign In", "✨ Create Account"])
        
        with tab1:
            email = st.text_input("Email", key="login_email", placeholder="you@example.com")
            password = st.text_input("Password", type="password", key="login_pass")
            
            if st.button("⚡ Enter", use_container_width=True, key="login_btn"):
                if email and password:
                    user_data, error = council.login_user(email, password)
                    if user_data:
                        st.session_state.user = user_data
                        # Save to localStorage for persistence across refreshes
                        st.markdown(f'''
                        <script>
                            saveSession("{user_data.get('token', '')}", "{user_data.get('id', '')}", "{user_data.get('email', '')}");
                        </script>
                        ''', unsafe_allow_html=True)
                        st.success("✅ Welcome back!")
                        st.rerun()
                    else:
                        st.error(f"❌ {error}")
                else:
                    st.warning("Enter email and password")
        
        with tab2:
            reg_email = st.text_input("Email", key="reg_email", placeholder="you@example.com")
            reg_pass = st.text_input("Password", type="password", key="reg_pass")
            reg_pass2 = st.text_input("Confirm", type="password", key="reg_pass2")
            
            if st.button("🚀 Create Account", use_container_width=True, key="reg_btn"):
                if reg_email and reg_pass:
                    if reg_pass != reg_pass2:
                        st.error("❌ Passwords don't match")
                    elif len(reg_pass) < 6:
                        st.error("❌ Min 6 characters")
                    else:
                        user_id, error = council.register_user(reg_email, reg_pass)
                        if user_id:
                            st.success("✅ Account created! Sign in now.")
                        else:
                            st.error(f"❌ {error}")
                else:
                    st.warning("Fill all fields")
    
    st.stop()

# ═══════════════════════════════════════════════════════════════════════════════
# MAIN APP
# ═══════════════════════════════════════════════════════════════════════════════

user_id = user["id"]

if "session_id" not in st.session_state:
    saved = params.get("session")
    if saved:
        st.session_state.session_id = saved
    else:
        st.session_state.session_id = council.create_session("New Chat", "Neon", user_id)
        st.query_params["session"] = st.session_state.session_id

if "artifact" not in st.session_state:
    st.session_state.artifact = "# 🧠 Ready to assist..."

if "screenshot" not in st.session_state:
    st.session_state.screenshot = None

# ═══════════════════════════════════════════════════════════════════════════════
# SIDEBAR
# ═══════════════════════════════════════════════════════════════════════════════

with st.sidebar:
    st.markdown(f"""
    <div style="text-align: center; padding: 25px 0;">
        <div style="font-size: 2.5em; margin-bottom: 10px;">🧠</div>
        <div style="
            font-size: 1.3em;
            background: linear-gradient(90deg, #ff1493, #ff6b35);
            -webkit-background-clip: text;
            -webkit-text-fill-color: transparent;
            font-weight: 700;
            letter-spacing: 1px;
        ">NEURAL COUNCIL</div>
        <div style="font-size: 0.8em; color: rgba(255,255,255,0.5); margin-top: 5px;">
            {user.get('email', 'User')[:25]}
        </div>
    </div>
    """, unsafe_allow_html=True)
    
    if st.button("🚪 Sign Out", use_container_width=True):
        logout()
    
    st.divider()
    
    if st.button("✨ New Chat", use_container_width=True):
        new_id = council.create_session("New Chat", "Neon", user_id)
        st.session_state.session_id = new_id
        st.query_params["session"] = new_id
        st.session_state.artifact = "# New chat..."
        council.reset_tokens()
        st.rerun()
    
    st.divider()
    
    st.markdown("### 🧠 The Council")
    for key, a in council.AGENTS.items():
        st.markdown(f"{a['avatar']} **{a['name']}**")
    
    st.divider()
    
    st.markdown("### 📂 Recent")
    try:
        for sess in council.get_sessions(user_id)[:8]:
            if sess['id'] == st.session_state.session_id:
                continue
            title = sess.get('title', 'Untitled')[:20]
            c1, c2 = st.columns([4, 1])
            with c1:
                if st.button(f"💬 {title}", key=f"s_{sess['id']}", use_container_width=True):
                    st.session_state.session_id = sess['id']
                    st.query_params["session"] = sess['id']
                    st.rerun()
            with c2:
                if st.button("×", key=f"d_{sess['id']}"):
                    council.delete_session(sess['id'])
                    st.rerun()
    except:
        st.caption("No history")
    
    st.divider()
    st.caption(f"⚡ {council.get_tokens_used():,} tokens")

# ═══════════════════════════════════════════════════════════════════════════════
# MAIN CONTENT
# ═══════════════════════════════════════════════════════════════════════════════

col1, col2 = st.columns([1.5, 1])

with col1:
    st.markdown("## 💬 Chat")
    
    with st.expander("📸 Screen Capture"):
        st.markdown(SCREEN_CAPTURE_JS, unsafe_allow_html=True)
        screenshot_b64 = st.text_area("Paste data:", key="ss_input", height=60, placeholder="Screenshot data...")
        if screenshot_b64 and screenshot_b64.startswith("data:"):
            st.session_state.screenshot = screenshot_b64
            st.success("✅ Ready!")
    
    # MULTI-FILE UPLOAD - ACCEPTS EVERYTHING
    uploaded_files = st.file_uploader(
        "📎 Attach Files", 
        accept_multiple_files=True,
        type=[
            # Code files
            'py', 'js', 'ts', 'jsx', 'tsx', 'java', 'c', 'cpp', 'h', 'cs', 'go', 'rs', 'rb', 'php', 'swift', 'kt',
            'scala', 'clj', 'ex', 'exs', 'erl', 'hs', 'ml', 'r', 'jl', 'lua', 'pl', 'pm', 'tcl', 'asm', 's',
            # Web files  
            'html', 'htm', 'css', 'scss', 'sass', 'less', 'vue', 'svelte', 'jsx', 'tsx', 'astro',
            # Data files
            'json', 'xml', 'yaml', 'yml', 'toml', 'csv', 'tsv', 'sql', 'graphql', 'prisma',
            # Documents
            'txt', 'md', 'rst', 'pdf', 'doc', 'docx', 'odt', 'rtf',
            # Spreadsheets
            'xlsx', 'xls', 'ods', 'numbers',
            # Presentations
            'pptx', 'ppt', 'odp', 'key',
            # Config files
            'env', 'ini', 'cfg', 'conf', 'sh', 'bash', 'zsh', 'fish', 'bat', 'ps1', 'cmd',
            # Images
            'png', 'jpg', 'jpeg', 'gif', 'webp', 'svg', 'bmp', 'ico', 'tiff', 'tif', 'heic', 'heif',
            # Videos
            'mp4', 'webm', 'mov', 'avi', 'mkv', 'flv', 'wmv', 'm4v', '3gp',
            # Audio (for transcription)
            'mp3', 'wav', 'm4a', 'ogg', 'flac', 'aac', 'wma', 'aiff',
            # Archives
            'zip', 'tar', 'gz', 'rar', '7z', 'bz2', 'xz',
            # Misc
            'log', 'gitignore', 'dockerfile', 'makefile', 'license', 'readme',
            # Jupyter
            'ipynb'
        ]
    )
    
    try:
        history = council.get_history(st.session_state.session_id)
        if history:
            # Group messages by conversation turns
            i = 0
            while i < len(history):
                msg = history[i]
                agent_name = msg.get("agent_name", "")
                role = msg["role"]
                
                # User messages always shown
                if role == "user":
                    with st.chat_message("user", avatar="👤"):
                        smart_display_content(msg["content"], is_user=True)  # Smart display for files
                    i += 1
                    continue
                
                # Check if this is a final answer
                is_final = any(x in str(agent_name) for x in ["Emperor", "Sage-Approved", "(Final)", "(Fallback)"])
                
                if is_final:
                    # Final answer - show prominently
                    avatar = "👑" if "Emperor" in str(agent_name) else "⚔️"
                    with st.chat_message("assistant", avatar=avatar):
                        st.markdown(f'<span class="agent-badge agent-emperor">✨ {agent_name}</span>', unsafe_allow_html=True)
                        render_message_with_images(msg["content"])
                        if "```" in str(msg["content"]):
                            try:
                                st.session_state.artifact = msg["content"].split("```")[1].split("\n", 1)[-1].strip()
                            except:
                                pass
                    i += 1
                else:
                    # Collect consecutive intermediate responses
                    intermediate = []
                    while i < len(history):
                        m = history[i]
                        a = m.get("agent_name", "")
                        if m["role"] == "user" or any(x in str(a) for x in ["Emperor", "Sage-Approved", "(Final)", "(Fallback)"]):
                            break
                        intermediate.append(m)
                        i += 1
                    
                    # Show intermediate responses in collapsed expander
                    if intermediate:
                        with st.expander(f"🔍 Council Deliberation ({len(intermediate)} steps)", expanded=False):
                            for m in intermediate:
                                a = m.get("agent_name", "")
                                cls = "strategist" if "Strategist" in a else "executor" if "Executor" in a else "sage" if "Sage" in a else ""
                                st.markdown(f'<span class="agent-badge agent-{cls}">{a}</span>', unsafe_allow_html=True)
                                st.markdown(m["content"])  # FULL content, no truncation
                                st.divider()
        else:
            st.info("✨ Start a conversation...")
    except Exception as e:
        st.error(f"⚠️ Error loading chat history: {str(e)}")
    
    # Input
    user_input = st.chat_input("Ask anything... (image: prompt, video: prompt, search: query)")
    
    if user_input:
        # Handle MULTIPLE uploaded files
        uploaded_images_b64 = []
        
        if uploaded_files:
            for uploaded in uploaded_files:
                try:
                    file_ext = uploaded.name.split('.')[-1].lower() if '.' in uploaded.name else ''
                    
                    if uploaded.type == "application/pdf":
                        from PyPDF2 import PdfReader
                        uploaded.seek(0)
                        file_text = "\n".join([p.extract_text() or "" for p in PdfReader(uploaded).pages])
                        user_input += f"\n\n[FILE: {uploaded.name}]\n```\n{file_text}\n```"
                    
                    elif uploaded.type and uploaded.type.startswith("image/"):
                        # Encode image as base64 so AI can SEE it
                        import base64
                        uploaded.seek(0)
                        image_bytes = uploaded.read()
                        img_b64 = f"data:{uploaded.type};base64,{base64.b64encode(image_bytes).decode()}"
                        uploaded_images_b64.append(img_b64)
                        user_input += f"\n\n[IMAGE ATTACHED: {uploaded.name}]"
                    
                    elif file_ext in ['mp4', 'webm', 'mov', 'avi', 'mkv']:
                        # VIDEO: Extract key frames for AI analysis
                        import tempfile
                        import subprocess
                        
                        user_input += f"\n\n[VIDEO: {uploaded.name} - {uploaded.size:,} bytes]"
                        
                        try:
                            # Save video temporarily
                            with tempfile.NamedTemporaryFile(suffix=f'.{file_ext}', delete=False) as tmp:
                                uploaded.seek(0)
                                tmp.write(uploaded.read())
                                tmp_path = tmp.name
                            
                            # Extract frames using ffmpeg (if available)
                            import shutil
                            if shutil.which('ffmpeg'):
                                # Extract first frame
                                frame_path = tmp_path + '_frame.jpg'
                                subprocess.run([
                                    'ffmpeg', '-i', tmp_path, '-vframes', '1', '-q:v', '2', frame_path
                                ], capture_output=True, timeout=30)
                                
                                # If frame extracted, add to images
                                import os
                                if os.path.exists(frame_path):
                                    with open(frame_path, 'rb') as f:
                                        frame_bytes = f.read()
                                    frame_b64 = f"data:image/jpeg;base64,{base64.b64encode(frame_bytes).decode()}"
                                    uploaded_images_b64.append(frame_b64)
                                    user_input += " [Frame extracted for analysis]"
                                    os.remove(frame_path)
                                
                                os.remove(tmp_path)
                            else:
                                user_input += " [Note: ffmpeg not installed - cannot extract frames]"
                        except Exception as ve:
                            user_input += f" [Frame extraction failed: {str(ve)[:50]}]"
                    
                    elif file_ext in ['mp3', 'wav', 'm4a', 'ogg', 'flac', 'aac', 'wma', 'aiff']:
                        # AUDIO: Note for transcription
                        user_input += f"\n\n[AUDIO: {uploaded.name} - {uploaded.size:,} bytes]"
                        user_input += "\n(Audio transcription: Use 'transcribe this audio' command if you want me to process it)"
                    
                    elif file_ext in ['xlsx', 'xls']:
                        # EXCEL: Parse spreadsheet
                        try:
                            import pandas as pd
                            uploaded.seek(0)
                            df = pd.read_excel(uploaded, engine='openpyxl' if file_ext == 'xlsx' else 'xlrd')
                            user_input += f"\n\n[EXCEL: {uploaded.name}]\n"
                            user_input += f"Columns: {list(df.columns)}\n"
                            user_input += f"Rows: {len(df)}\n"
                            user_input += f"```\n{df.to_string()}\n```"
                        except Exception as ee:
                            user_input += f"\n\n[EXCEL: {uploaded.name} - Could not parse: {str(ee)[:50]}]"
                    
                    elif file_ext == 'csv':
                        # CSV: Parse and show
                        try:
                            import pandas as pd
                            uploaded.seek(0)
                            df = pd.read_csv(uploaded)
                            user_input += f"\n\n[CSV: {uploaded.name}]\n"
                            user_input += f"Columns: {list(df.columns)}\n"
                            user_input += f"Rows: {len(df)}\n"
                            user_input += f"```\n{df.to_string()}\n```"
                        except Exception as ce:
                            user_input += f"\n\n[CSV: {uploaded.name} - Could not parse: {str(ce)[:50]}]"
                    
                    elif file_ext in ['docx']:
                        # WORD: Extract text
                        try:
                            from docx import Document
                            uploaded.seek(0)
                            doc = Document(uploaded)
                            doc_text = "\n".join([p.text for p in doc.paragraphs])
                            user_input += f"\n\n[WORD: {uploaded.name}]\n```\n{doc_text}\n```"
                        except Exception as we:
                            user_input += f"\n\n[WORD: {uploaded.name} - Could not parse: {str(we)[:50]}]"
                    
                    elif file_ext in ['pptx']:
                        # POWERPOINT: Extract text from slides
                        try:
                            from pptx import Presentation
                            uploaded.seek(0)
                            ppt = Presentation(uploaded)
                            slides_text = []
                            for i, slide in enumerate(ppt.slides[:20]):
                                slide_text = f"--- Slide {i+1} ---\n"
                                for shape in slide.shapes:
                                    if hasattr(shape, "text"):
                                        slide_text += shape.text + "\n"
                                slides_text.append(slide_text)
                            user_input += f"\n\n[POWERPOINT: {uploaded.name}]\n```\n{chr(10).join(slides_text)}\n```"
                        except Exception as pe:
                            user_input += f"\n\n[POWERPOINT: {uploaded.name} - Could not parse: {str(pe)[:50]}]"
                    
                    elif file_ext == 'ipynb':
                        # JUPYTER: Extract code cells
                        try:
                            uploaded.seek(0)
                            nb = json.loads(uploaded.read().decode('utf-8'))
                            cells = []
                            for cell in nb.get('cells', [])[:30]:
                                if cell.get('cell_type') == 'code':
                                    cells.append("```python\n" + "".join(cell.get('source', [])) + "\n```")
                                elif cell.get('cell_type') == 'markdown':
                                    cells.append("".join(cell.get('source', [])))
                            user_input += f"\n\n[JUPYTER: {uploaded.name}]\n" + "\n\n".join(cells)
                        except Exception as je:
                            user_input += f"\n\n[JUPYTER: {uploaded.name} - Could not parse: {str(je)[:50]}]"
                    
                    elif file_ext in ['zip', 'tar', 'gz', 'rar', '7z']:
                        # For archives, just note them
                        user_input += f"\n\n[ARCHIVE: {uploaded.name} - {uploaded.size:,} bytes]"
                    
                    else:
                        # Read as text - larger limit for code files
                        uploaded.seek(0)
                        try:
                            file_text = uploaded.read().decode('utf-8', errors='ignore')
                        except:
                            file_text = str(uploaded.read())
                        
                        # Detect language for syntax highlighting
                        lang_map = {
                            'py': 'python', 'js': 'javascript', 'ts': 'typescript', 'jsx': 'jsx', 'tsx': 'tsx',
                            'java': 'java', 'c': 'c', 'cpp': 'cpp', 'h': 'c', 'cs': 'csharp', 'go': 'go',
                            'rs': 'rust', 'rb': 'ruby', 'php': 'php', 'swift': 'swift', 'kt': 'kotlin',
                            'html': 'html', 'css': 'css', 'scss': 'scss', 'json': 'json', 'xml': 'xml',
                            'yaml': 'yaml', 'yml': 'yaml', 'sql': 'sql', 'md': 'markdown', 'sh': 'bash'
                        }
                        lang = lang_map.get(file_ext, '')
                        
                        # NO LIMIT - AI sees complete files!
                        user_input += f"\n\n[FILE: {uploaded.name}]\n```{lang}\n{file_text}\n```"
                            
                except Exception as e:
                    user_input += f"\n\n[FILE ERROR: {uploaded.name} - {str(e)}]"
        
        # CACHE FILES: Store files for efficient follow-up messages
        # Extract file info and cache it
        import re
        file_matches = re.findall(r'\[FILE: ([^\]]+)\]\n```[^\n]*\n([\s\S]*?)```', user_input)
        if file_matches:
            files_to_cache = [{"name": name, "content": content} for name, content in file_matches]
            council.cache_files(st.session_state.session_id, files_to_cache)
        
        with st.chat_message("user", avatar="👤"):
            smart_display_content(user_input, is_user=True)  # Smart display collapses files
            # Show image thumbnails if any
            if uploaded_images_b64:
                st.caption(f"📷 {len(uploaded_images_b64)} image(s) attached")
        
        # Use ALL uploaded images OR screenshot (combine into single for now, vision API takes first)
        # Note: Claude's vision API can handle multiple images - we join them
        if uploaded_images_b64:
            # For multiple images, we'll use the first one for now (API limitation)
            # but add note about others
            screenshot = uploaded_images_b64[0]
            if len(uploaded_images_b64) > 1:
                user_input += f"\n\n[Note: {len(uploaded_images_b64)} images attached, analyzing first image]"
        else:
            screenshot = st.session_state.get("screenshot")
        st.session_state.screenshot = None
        
        with st.status("⚡ Council processing...", expanded=True) as status:
            try:
                final_answer = None
                final_agent = None
                intermediate_responses = []
                
                for agent, content, msg_type in council.run_council("Neon", user_input, st.session_state.session_id, user_id, screenshot):
                    if msg_type == "system":
                        st.markdown(f"🔔 {content}")
                    elif msg_type == "image":
                        st.image(content, width=500)
                    elif msg_type == "video":
                        st.video(content)
                    else:
                        # Check if this is the final answer (Emperor or Sage-Approved or Fallback)
                        is_final = any(x in agent for x in ["Emperor", "Sage-Approved", "(Final)", "(Fallback)"])
                        
                        if is_final:
                            final_answer = content
                            final_agent = agent
                        else:
                            # Store intermediate response for expander
                            avatar = next((a["avatar"] for a in council.AGENTS.values() if a["name"] in str(agent)), "🤖")
                            intermediate_responses.append((agent, content, avatar))
                        
                        if "```" in str(content):
                            try:
                                st.session_state.artifact = content.split("```")[1].split("\n", 1)[-1].strip()
                            except:
                                pass
                
                status.update(label="✅ Complete!", state="complete")
            except Exception as e:
                status.update(label=f"❌ {str(e)}", state="error")
        
        # Display intermediate responses in expanders
        if intermediate_responses:
            with st.expander("🔍 View Council Deliberation", expanded=False):
                for agent, content, avatar in intermediate_responses:
                    cls = "strategist" if "Strategist" in agent else "executor" if "Executor" in agent else "sage" if "Sage" in agent else ""
                    st.markdown(f'<span class="agent-badge agent-{cls}">{agent}</span>', unsafe_allow_html=True)
                    smart_display_content(content)  # Smart display for long responses
                    st.divider()
        
        # Display final answer prominently
        if final_answer:
            avatar = "👑" if "Emperor" in str(final_agent) else "⚔️"
            with st.chat_message("assistant", avatar=avatar):
                st.markdown(f'<span class="agent-badge agent-emperor">✨ {final_agent}</span>', unsafe_allow_html=True)
                render_message_with_images(final_answer)
        
        st.rerun()

with col2:
    st.markdown("## 📄 Artifacts")
    
    code = st.session_state.artifact
    
    # Smart language detection
    lang = "text"
    ext = "txt"
    if code.strip().startswith(("{", "[")):
        lang, ext = "json", "json"
    elif "<html" in code.lower() or "<!doctype" in code.lower():
        lang, ext = "html", "html"
    elif "function " in code or "const " in code or "=>" in code:
        lang, ext = "javascript", "js"
    elif "def " in code or "import " in code or "class " in code:
        lang, ext = "python", "py"
    elif "package " in code or "public class" in code:
        lang, ext = "java", "java"
    elif code.strip().startswith("/*") or "#include" in code:
        lang, ext = "c", "c"
    elif "<" in code and ">" in code and "/" in code:
        lang, ext = "xml", "xml"
    elif "---" in code and ":" in code:
        lang, ext = "yaml", "yaml"
    elif "SELECT " in code.upper() or "INSERT " in code.upper():
        lang, ext = "sql", "sql"
    elif code.strip().startswith("#"):
        lang, ext = "markdown", "md"
    elif code.strip().startswith("@") or "color:" in code:
        lang, ext = "css", "css"
    
    st.code(code, language=lang, line_numbers=True)
    
    # Custom filename input
    default_name = f"artifact.{ext}"
    filename = st.text_input("📝 Filename", value=default_name, key="artifact_filename", placeholder="myfile.py")
    
    c1, c2 = st.columns(2)
    with c1:
        st.download_button("💾 Download", code, filename, use_container_width=True)
    with c2:
        if st.button("🗑️ Clear", use_container_width=True):
            st.session_state.artifact = "# Cleared"
            st.rerun()
    
    st.divider()
    
    st.markdown("### 🧠 The Models")
    st.markdown("""
    | Tier | Agent | Model |
    |------|-------|-------|
    | 👑 | Emperor | Claude Opus 4.5 |
    | 🎯 | Strategist | Claude Sonnet 4.5 |
    | ⚔️ | Executor | GPT-5.2 |
    | 📿 | Sage | DeepSeek V3.2 |
    """)
    
    with st.expander("💡 Commands"):
        st.markdown("""
        - `image: [prompt]` → Generate image
        - `video: [prompt]` → Generate video
        - `search: [query]` → Web search
        - Natural: "create an image of..."
        - Natural: "make a video of..."
        """)
